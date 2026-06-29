# ============================================================
# 檔案名稱：parsers.py
# 中文名稱：文件文字萃取模組
# 功能說明：從 pptx / docx / pdf / txt 中萃取純文字，供 Ingest LLM 使用
# 所屬模組：job/
# 建立日期：2026-06-28
# 修改日期：2026-06-28
# 開發者　：Blackjtsai
# ============================================================

import os
from pathlib import Path


def extract_pptx(file_path: str) -> str:
    """從 .pptx 萃取所有投影片文字（含標題、內容、備忘錄）。

    Args:
        file_path: pptx 檔案絕對路徑

    Returns:
        純文字（每張投影片以分隔線區隔）
    """
    from pptx import Presentation

    prs = Presentation(file_path)
    slides_text: list[str] = []

    for i, slide in enumerate(prs.slides, start=1):
        parts: list[str] = [f"--- 投影片 {i} ---"]
        for shape in slide.shapes:
            if shape.has_text_frame:
                for para in shape.text_frame.paragraphs:
                    line = para.text.strip()
                    if line:
                        parts.append(line)
        # 備忘錄
        if slide.has_notes_slide:
            notes = slide.notes_slide.notes_text_frame.text.strip()
            if notes:
                parts.append(f"[備忘錄] {notes}")
        slides_text.append("\n".join(parts))

    return "\n\n".join(slides_text)


def extract_docx(file_path: str) -> str:
    """從 .docx 萃取所有段落文字。

    Args:
        file_path: docx 檔案絕對路徑

    Returns:
        純文字
    """
    from docx import Document

    doc = Document(file_path)
    paragraphs = [p.text.strip() for p in doc.paragraphs if p.text.strip()]
    return "\n\n".join(paragraphs)


def extract_pdf(file_path: str) -> str:
    """從 .pdf 萃取所有頁面文字。

    Args:
        file_path: pdf 檔案絕對路徑

    Returns:
        純文字
    """
    from pypdf import PdfReader

    reader = PdfReader(file_path)
    pages_text: list[str] = []
    for i, page in enumerate(reader.pages, start=1):
        text = page.extract_text() or ""
        if text.strip():
            pages_text.append(f"--- 第 {i} 頁 ---\n{text.strip()}")

    return "\n\n".join(pages_text)


def extract_txt(file_path: str) -> str:
    """從純文字檔讀取內容（自動偵測 encoding）。

    Args:
        file_path: txt 檔案絕對路徑

    Returns:
        純文字
    """
    for encoding in ("utf-8", "utf-8-sig", "big5", "gbk", "cp950"):
        try:
            with open(file_path, encoding=encoding) as f:
                return f.read()
        except (UnicodeDecodeError, LookupError):
            continue

    # 最後退路：忽略錯誤字元
    with open(file_path, encoding="utf-8", errors="replace") as f:
        return f.read()


def extract_text(file_path: str) -> str:
    """根據副檔名自動選擇解析器，萃取純文字。

    Args:
        file_path: 任意支援格式的檔案絕對路徑

    Returns:
        純文字

    Raises:
        ValueError: 不支援的檔案格式
    """
    suffix = Path(file_path).suffix.lower()
    dispatchers = {
        ".pptx": extract_pptx,
        ".ppt":  extract_pptx,   # python-pptx 支援舊格式（有限）
        ".docx": extract_docx,
        ".doc":  extract_docx,
        ".pdf":  extract_pdf,
        ".txt":  extract_txt,
        ".md":   extract_txt,
    }

    parser = dispatchers.get(suffix)
    if parser is None:
        raise ValueError(f"不支援的檔案格式：{suffix}（支援：pptx/docx/pdf/txt/md）")

    return parser(file_path)


SUPPORTED_EXTENSIONS = {".pptx", ".ppt", ".docx", ".doc", ".pdf", ".txt", ".md"}
